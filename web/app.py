"""
app.py — Web Server & Feature Layer
=======================================
This is where ALL your features live.
agent.py stays the same — only this file changes per feature.

Current Routes:
  GET  /          → Renders the browser UI (index.html)
  POST /hello     → Hello World feature (greets the user)

How to add a new feature:
  1. Add a new HTML section in index.html
  2. Add a new @app.route() here
  3. Build a prompt string from the user's input
  4. Call ask_agent(prompt) and return the result

Flow:
  Browser → app.py (builds prompt) → agent.py (calls Claude) → Browser
"""

import base64
import os
import re
import subprocess
import sys
import tempfile

# Make sure Python can find the agent/ folder from anywhere
sys.path.insert(0, os.path.abspath(os.path.dirname(__file__) + "/.."))

from flask import Flask, render_template, request, jsonify
from dotenv import load_dotenv
from agent.agent import ask_agent

# Directory where generated prompts are saved as .md files
PROMPTS_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "generated-prompts"))

# Voice reference samples for Chatterbox TTS cloning
_VOICE_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
VOICE_SAMPLES = {
    "sekhar":  os.path.join(_VOICE_DIR, "sekhar_voice.wav"),
    "abhinav": os.path.join(_VOICE_DIR, "abhinav_voice.wav"),
    "akshay":  os.path.join(_VOICE_DIR, "akshay_voice.wav"),
}

# Chatterbox model — loaded once on first request
_chatterbox = None

def get_chatterbox():
    global _chatterbox
    if _chatterbox is None:
        from chatterbox.tts import ChatterboxTTS
        _chatterbox = ChatterboxTTS.from_pretrained(device="cpu")
    return _chatterbox


def git_push_prompt(filepath: str, filename: str) -> None:
    """Commit the newly saved prompt file and push to github.com/badugusekhar/claude-code-prompts."""
    try:
        subprocess.run(["git", "-C", PROMPTS_DIR, "add", filename], check=True, timeout=15)
        subprocess.run(
            ["git", "-C", PROMPTS_DIR, "commit", "-m", f"add: {filename}"],
            check=True, timeout=15,
        )
        subprocess.run(["git", "-C", PROMPTS_DIR, "push"], check=True, timeout=30)
    except Exception:
        pass  # push failure is non-fatal — prompt is still saved locally


def slugify(text: str) -> str:
    """Turn a task description into a kebab-case filename slug (max 50 chars)."""
    slug = text.lower()
    slug = re.sub(r"[^a-z0-9\s-]", "", slug)
    slug = re.sub(r"\s+", "-", slug.strip())
    slug = re.sub(r"-+", "-", slug)
    return slug[:50].rstrip("-")

# Load ANTHROPIC_API_KEY from .env file
load_dotenv()

# Resolve templates/static relative to this file so Flask works
# whether the script is run from the project root or the web folder.
_base_dir = os.path.dirname(__file__)
app = Flask(
    __name__,
    template_folder=os.path.join(_base_dir, "templates"),
    static_folder=os.path.join(_base_dir, "static"),
)


# ─────────────────────────────────────────────
# ROUTE 1: Home Page
# ─────────────────────────────────────────────
@app.route("/")
def index():
    """Serve the main browser UI."""
    return render_template("index.html")


# ─────────────────────────────────────────────
# ROUTE 2: Hello World Feature
# ─────────────────────────────────────────────
@app.route("/hello", methods=["POST"])
def hello():
    """
    Feature: Greet the user by name.

    Browser sends: { "name": "Sekhar" }
    We build:      "Say a friendly hello to Sekhar in one sentence."
    Claude returns: "Hello, Sekhar! ..."
    Browser shows:  the response
    """
    data = request.get_json()
    name = data.get("name", "World").strip()

    # Build a natural language prompt — this is the ONLY thing you change per feature
    prompt = (
        f"Say a warm and friendly hello to {name} in exactly one sentence. "
        f"Keep it simple and cheerful."
    )

    response = ask_agent(prompt)
    return jsonify({"message": response})


# ─────────────────────────────────────────────
# ROUTE 3: Generate Prompt Feature
# ─────────────────────────────────────────────
@app.route("/generate-prompt", methods=["POST"])
def generate_prompt():
    """
    Feature: Expand a short task description into a detailed Claude Code prompt.

    Browser sends: { "task": "add a dark mode toggle" }
    Returns:       { "prompt": "...", "filename": "add-dark-mode-toggle.md" }
    Also saves the prompt as generated-prompts/<slug>.md
    """
    data = request.get_json()
    task = data.get("task", "").strip()
    if not task:
        return jsonify({"error": "No task provided"}), 400

    meta_prompt = f"""You are a Claude Code prompt engineer. Convert the task below into a concise, high-signal prompt — 150 to 250 words maximum. Claude Code is already skilled at exploring codebases; it does not need hand-holding. Give it intent and boundaries, not a tutorial.

Task: "{task}"

Output ONLY the prompt, using this exact structure. Keep each section tight:

**Task**
One or two sentences. What to build or fix and the expected outcome. Be specific, not generic.

**Constraints**
3 bullets max. Only include non-obvious guardrails — things Claude would not naturally assume. Skip obvious ones like "don't break tests".

**Done when**
3 bullets max. Each must be independently verifiable by running or clicking something. No vague conditions.

**Start here**
A short list of files most relevant to this task. Infer likely filenames from the task description and common project layouts.

Rules:
- Total output must be under 250 words
- No fluff, no re-stating the obvious, no step-by-step approach (Claude figures that out)
- No preamble or explanation outside the four sections above"""

    response = ask_agent(meta_prompt)

    # Save to generated-prompts/<slug>.md
    slug = slugify(task)
    filename = f"{slug}.md"
    os.makedirs(PROMPTS_DIR, exist_ok=True)
    filepath = os.path.join(PROMPTS_DIR, filename)
    # Avoid clobbering — append a counter if the file already exists
    counter = 1
    while os.path.exists(filepath):
        filepath = os.path.join(PROMPTS_DIR, f"{slug}-{counter}.md")
        filename = f"{slug}-{counter}.md"
        counter += 1
    with open(filepath, "w", encoding="utf-8") as f:
        f.write(f"# {task}\n\n")
        f.write(response)
        f.write("\n")

    git_push_prompt(filepath, filename)

    return jsonify({"prompt": response, "filename": filename})


# ─────────────────────────────────────────────
# ROUTE 4: Scrum Status Voice
# ─────────────────────────────────────────────
@app.route("/scrum-voice", methods=["POST"])
def scrum_voice():
    """
    Feature: Generate a voice recording of a daily scrum status update.

    Browser sends: { "today": "...", "tomorrow": "...", "blockers": "..." }
    Returns:       { "script": "...", "audio": "data:audio/wav;base64,..." }

    Step 1 — Claude rewrites bullet points as natural spoken language.
    Step 2 — Coqui XTTS-v2 synthesises audio using sekhar_voice.m4a as the
              voice reference (runs fully locally, no API key needed).
    """
    data     = request.get_json()
    today    = data.get("today", "").strip()
    tomorrow = data.get("tomorrow", "").strip()
    blockers = data.get("blockers", "").strip() or "None"
    voice    = data.get("voice", "sekhar").strip().lower()

    if not today and not tomorrow:
        return jsonify({"error": "Please fill in at least Today or Tomorrow."}), 400

    voice_sample = VOICE_SAMPLES.get(voice, VOICE_SAMPLES["sekhar"])
    if not os.path.exists(voice_sample):
        return jsonify({"error": f"{voice.capitalize()}'s voice file not found. Please copy {voice}_voice.wav to the project root folder."}), 400

    rewrite_prompt = f"""Rewrite this daily scrum status as natural, conversational spoken English for a team standup — the kind you'd say out loud, not read from a list. Use contractions and smooth transitions. Keep it under 45 seconds when spoken aloud. Output ONLY the spoken script, no labels or preamble.

Today: {today}
Tomorrow: {tomorrow}
Blockers: {blockers}"""

    script = ask_agent(rewrite_prompt)

    try:
        import torchaudio
        import numpy as np
        import librosa
        model = get_chatterbox()
        wav = model.generate(script, audio_prompt_path=voice_sample, exaggeration=0.0, cfg_weight=0.95)
        # Slow down to 80% speed to match natural speaking pace (no pitch change)
        wav_np = wav.squeeze().cpu().numpy()
        wav_slow = librosa.effects.time_stretch(wav_np, rate=0.80)
        wav_out = np.expand_dims(wav_slow, axis=0)
        import torch
        wav_tensor = torch.from_numpy(wav_out)
        tmp_path = tempfile.mktemp(suffix=".wav")
        torchaudio.save(tmp_path, wav_tensor, model.sr)
        with open(tmp_path, "rb") as f:
            audio_bytes = f.read()
        os.unlink(tmp_path)
    except Exception as e:
        return jsonify({"error": f"TTS generation failed: {str(e)}"}), 500

    audio_b64 = base64.b64encode(audio_bytes).decode("utf-8")
    return jsonify({"script": script, "audio": f"data:audio/wav;base64,{audio_b64}"})


# ─────────────────────────────────────────────
# ROUTE 5: Weekly Status Report PPT Generator
# ─────────────────────────────────────────────

WEEKLY_TEMPLATE = r"C:\Users\sekhar.a.badugu\OneDrive - Accenture\T-Mobile\weekly_Status_Report\2026\T-Mobile_Monitoring Dev Team_Project Status_2026-07-17.pptx"

@app.route("/weekly-report", methods=["POST"])
def weekly_report():
    """
    Feature: Generate a weekly T-Mobile status PowerPoint report.

    Browser sends: { "date": "2026-07-24", "key_updates": "...", "next_steps": "...", "blockers": "..." }
    Returns:       { "filename": "...", "file": "data:application/...;base64,..." }

    Copies the template PPTX, replaces the date on slide 1 and fills the
    three content rows on slide 2 using python-pptx. Claude polishes each
    section into professional bullet-point text first.
    """
    import json as _json
    import re as _re
    import shutil
    from pptx import Presentation
    from pptx.util import Pt

    data        = request.get_json()
    week_date   = data.get("date", "").strip()
    key_updates = data.get("key_updates", "").strip()
    next_steps  = data.get("next_steps", "").strip()
    blockers    = data.get("blockers", "").strip() or "None"

    if not week_date or not key_updates:
        return jsonify({"error": "Date and Key Updates are required."}), 400

    # Step 1: Claude polishes content into structured bullets
    polish_prompt = f"""You are writing a professional weekly status report for a software development team at T-Mobile.
Format each section as clean bullet points. Use sub-topic labels followed by bullet detail lines.
Keep language concise and professional. Do NOT add any extra commentary.

Key Updates:
{key_updates}

Next Steps:
{next_steps}

Blockers/Dependencies:
{blockers}

Return ONLY a valid JSON object with exactly these three keys: "key_updates", "next_steps", "blockers".
Each value is a list of strings. Sub-topic label lines end with a colon (e.g. "OneConsole:").
Bullet detail lines start with "• " (e.g. "• Completed integration testing").
Example format:
{{
  "key_updates": ["OneConsole:", "• Fixed pipeline issue", "PITSTOP:", "• Ready for deployment"],
  "next_steps": ["• Take new story if available", "• Continue integration testing"],
  "blockers": ["• MR reviews pending"]
}}"""

    raw = ask_agent(polish_prompt)

    match = _re.search(r'\{.*\}', raw, _re.DOTALL)
    try:
        sections = _json.loads(match.group(0)) if match else {}
    except Exception:
        sections = {}

    if not sections.get("key_updates"):
        sections = {
            "key_updates": [key_updates],
            "next_steps": [next_steps] if next_steps else ["• No updates"],
            "blockers": [blockers],
        }

    # Step 2: Copy template and write content with python-pptx
    tmp = tempfile.mktemp(suffix=".pptx")
    shutil.copy2(WEEKLY_TEMPLATE, tmp)
    prs = Presentation(tmp)

    # Slide 1: replace the date (3rd paragraph of the title text box)
    slide1 = prs.slides[0]
    for shape in slide1.shapes:
        if shape.has_text_frame and len(shape.text_frame.paragraphs) >= 3:
            para = shape.text_frame.paragraphs[2]
            for run in para.runs:
                run.text = week_date
            break

    # Slide 2 has 3 separate tables (Key Updates, Next Steps, Blockers)
    # Each table has 2 rows: row[0]=header, row[1]=content
    slide2 = prs.slides[1]
    tables = [s.table for s in slide2.shapes if s.has_table]

    def fill_cell(cell, lines):
        tf = cell.text_frame
        tf.clear()
        for i, line in enumerate(lines):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            is_label = line.rstrip().endswith(":") and not line.startswith("•")
            p.level = 0 if is_label else 1
            run = p.add_run()
            run.text = line.lstrip("• ")
            run.font.bold = is_label
            run.font.size = Pt(16) if is_label else Pt(14)

    if len(tables) >= 1:
        fill_cell(tables[0].rows[1].cells[0], sections.get("key_updates", [key_updates]))
    if len(tables) >= 2:
        fill_cell(tables[1].rows[1].cells[0], sections.get("next_steps", [next_steps or "• No updates"]))
    if len(tables) >= 3:
        fill_cell(tables[2].rows[1].cells[0], sections.get("blockers", [blockers]))

    prs.save(tmp)
    with open(tmp, "rb") as f:
        pptx_bytes = f.read()
    os.unlink(tmp)

    filename = f"T-Mobile_Monitoring Dev Team_Project Status_{week_date}.pptx"
    pptx_b64 = base64.b64encode(pptx_bytes).decode("utf-8")
    return jsonify({
        "filename": filename,
        "file": f"data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{pptx_b64}"
    })


# ─────────────────────────────────────────────
# ROUTE 6: Send Weekly Report via Outlook
# ─────────────────────────────────────────────
@app.route("/send-email", methods=["POST"])
def send_email():
    """
    Feature: Open a pre-filled Outlook draft for the weekly status report.

    Browser sends: { date, key_updates, next_steps, blockers, pptx_b64, filename }
    Returns:       { "ok": true }

    Claude writes a short email body summary, then win32com opens Outlook with
    To/CC/Subject/Body/Attachment pre-filled. The user reviews and clicks Send.
    """
    import win32com.client

    data        = request.get_json()
    week_date   = data.get("date", "").strip()
    key_updates = data.get("key_updates", "").strip()
    next_steps  = data.get("next_steps", "").strip()
    blockers    = data.get("blockers", "").strip() or "None"
    pptx_b64    = data.get("pptx_b64", "")
    filename    = data.get("filename", f"Status_{week_date}.pptx")

    if not pptx_b64:
        return jsonify({"error": "No PPTX data — generate the report first."}), 400

    # Step 1: Claude writes the email body
    summary_prompt = f"""Write a short professional email body for a weekly T-Mobile status report email.
The email is from Sekhar Badugu to Grace (project manager).

Key Updates this week:
{key_updates}

Next Steps:
{next_steps}

Blockers:
{blockers}

Format:
- Start with: "Hello Grace,"
- One sentence intro: "Please see attached for the weekly status update of the Monitoring Dev Team(OneConsole)."
- "To Summarize:" followed by 3-6 bullet points of the most important highlights
- End with: "Let me know if you have any questions or concerns."
- Do NOT include the sign-off — it will be added separately.
- Output ONLY the email body text, no subject line."""

    body_text = ask_agent(summary_prompt)

    sign_off = (
        "\n\nBest Regards,\n"
        "Sekhar Badugu (Sekhar)\n"
        "sekhar.a.badugu@accenture.com\n"
        "Mobile: +91 9740252232\n"
        "Upcoming PTO: None"
    )
    full_body = body_text.strip() + sign_off

    # Step 2: Write PPTX bytes to temp file for attachment
    pptx_bytes = base64.b64decode(pptx_b64.split(",", 1)[-1])
    tmp_pptx = os.path.join(tempfile.gettempdir(), filename)
    with open(tmp_pptx, "wb") as f:
        f.write(pptx_bytes)

    # Step 3: Open Outlook draft via COM automation
    try:
        outlook = win32com.client.Dispatch("Outlook.Application")
        mail = outlook.CreateItem(0)  # 0 = olMailItem
        mail.To = "grace.claflin1@t-mobile.com"
        mail.CC = (
            "Deshpande, Neelkanth; Harika, Tegbir; Rastogi, Nitin; "
            "Gadhave, Akshay; Agarwal, Abhinav; Stratton, Jesse"
        )
        mail.Subject = f"T-Mobile | Monitoring Dev Team | Weekly Status Update | {week_date}"
        mail.Body = full_body
        mail.Attachments.Add(tmp_pptx)
        mail.Display()  # opens draft — user clicks Send
    except Exception as e:
        return jsonify({"error": f"Outlook error: {str(e)}"}), 500

    return jsonify({"ok": True})


# ─────────────────────────────────────────────
# FUTURE FEATURES — Add new routes below here
# ─────────────────────────────────────────────

# Example skeleton for next feature (Calculator):
#
# @app.route("/calculate", methods=["POST"])
# def calculate():
#     data = request.get_json()
#     num1 = data.get("num1")
#     num2 = data.get("num2")
#     operation = data.get("operation", "plus")
#     prompt = f"Calculate {num1} {operation} {num2}. Reply with just the number."
#     response = ask_agent(prompt)
#     return jsonify({"result": response})


if __name__ == "__main__":
    print("=" * 50)
    print("  Claude Agents Playground — Starting Server")
    print("  Open your browser at: http://localhost:5000")
    print("=" * 50)
    app.run(debug=True, port=5000)
